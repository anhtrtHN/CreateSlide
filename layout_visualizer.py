
from pptx import Presentation
from pptx.util import Length
import sys

def draw_slide_ascii(slide, index):
    # Slide Dimensions (16:9)
    # Width ~ 33.867 cm
    # Height ~ 19.05 cm
    
    # Scale: 1 character horiz = 0.5 cm
    # Scale: 1 character vert = 0.5 cm
    # Canvas Size: ~68 chars wide x ~38 chars high
    
    W_scale = 0.5
    H_scale = 0.5
    
    width_chars = int(33.867 / W_scale)
    height_chars = int(19.05 / H_scale)
    
    canvas = [[' ' for _ in range(width_chars)] for _ in range(height_chars)]
    
    def draw_rect(top_cm, left_cm, width_cm, height_cm, char, fill=False):
        t = int(top_cm / H_scale)
        l = int(left_cm / W_scale)
        w = int(width_cm / W_scale)
        h = int(height_cm / H_scale)
        
        for r in range(t, min(t+h, height_chars)):
            for c in range(l, min(l+w, width_chars)):
                if r == t or r == t+h-1 or c == l or c == l+w-1:
                    canvas[r][c] = char
                elif fill:
                    canvas[r][c] = '.'
    
    # Draw Border
    draw_rect(0, 0, 33.8, 19.0, '#')
    
    title_shape = slide.shapes.title
    
    # Find Body
    body_shape = None
    # Quick heuristics to find the body used in engine
    # We will just iterate and find the one that ISN'T title
    for shape in slide.placeholders:
        if shape.element is not title_shape.element:
            body_shape = shape
            break
            
    # Draw Title (T)
    if title_shape:
        draw_rect(title_shape.top.cm, title_shape.left.cm, title_shape.width.cm, title_shape.height.cm, 'T')
        
    # Draw Body (B)
    if body_shape:
         draw_rect(body_shape.top.cm, body_shape.left.cm, body_shape.width.cm, body_shape.height.cm, 'B', fill=True)
         
    print(f"\n--- SLIDE {index + 1}: {title_shape.text[:30]}... ---")
    if title_shape:
        print(f"Title: Top={title_shape.top.cm:.2f}cm, H={title_shape.height.cm:.2f}cm")
    if body_shape:
        gap = body_shape.top.cm - (title_shape.top.cm + title_shape.height.cm) if title_shape else 0
        print(f"Body : Top={body_shape.top.cm:.2f}cm, H={body_shape.height.cm:.2f}cm")
        print(f"GAP  : {gap:.2f} cm")
        
    print("-" * width_chars)
    for row in canvas:
        print("".join(row))
    print("-" * width_chars)

def main():
    try:
        prs = Presentation("stress_test.pptx")
        for i, slide in enumerate(prs.slides):
            draw_slide_ascii(slide, i)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    main()
