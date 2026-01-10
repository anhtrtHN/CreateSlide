
import json
import slide_engine
from pptx import Presentation
from pptx.util import Cm

def run_reproduction():
    # 1. Create Data with Long Title
    data = {
        "title": "Main Presentation",
        "slides": [
            {
                "title": "This is a very long title that should definitely wrap to at least two lines because it has many many words in it to test the layout engine capabilities",
                "content": ["Bullet 1", "Bullet 2"]
            },
             {
                "title": "Short Title",
                "content": ["Bullet 1", "Bullet 2"]
            }
        ]
    }
    
    # 2. Generate PPTX
    print("Generating PPTX...")
    pptx_bytes = slide_engine.create_pptx(data)
    
    # 3. Save to disk for manual inspection if needed
    filename = "reproduce_layout.pptx"
    with open(filename, "wb") as f:
        f.write(pptx_bytes.getbuffer())
        
    # 4. Inspect Properties
    prs = Presentation(filename)
    
    print("\n--- INSPECTION ---")
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title
        print(f"\nSlide {i}: Title Text: '{title.text}'")
        if title:
            print(f"Title Height: {title.height.cm:.2f} cm")
            print(f"Title Top: {title.top.cm:.2f} cm")
            bottom_cm = title.top.cm + title.height.cm
            print(f"Title Bottom: {bottom_cm:.2f} cm")
            
        # Find body
        body = None
        for shape in slide.placeholders:
            if shape.element is not title.element:
                body = shape
                break
        
        if body:
            print(f"Body Top: {body.top.cm:.2f} cm")
            gap = body.top - (title.top + title.height)
            print(f"Gap: {gap / 360000:.2f} cm")
        else:
            print("No body found")

if __name__ == "__main__":
    run_reproduction()
