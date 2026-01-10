import io
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def create_pptx(json_data: Dict[str, Any], template_pptx_bytes: bytes | None = None) -> io.BytesIO:
    """
    Generates a PowerPoint presentation from JSON data.
    Returns: BytesIO object of the .pptx file.
    """
    if template_pptx_bytes:
        prs = Presentation(io.BytesIO(template_pptx_bytes))
        # Clear existing slides from the template explicitly and safely
        # We need to remove the relationship (rId) to avoid corruption
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        
        for s in slides:
             rId = s.rId
             prs.part.drop_rel(rId) # Critical: Remove the relationship
             xml_slides.remove(s)   # Remove the slide entry
            
        print(f"Template loaded and cleared. Remaining slides: {len(prs.slides)}")
        
    else:
        prs = Presentation() # Uses default template

    # Helper to find usable layouts
    def get_layout(prs, preferred_index, needs_body=False):
        # Try preferred index first
        if preferred_index < len(prs.slide_layouts):
            layout = prs.slide_layouts[preferred_index]
            if not needs_body:
                return layout
            # Check if it has a body placeholder (usually idx 1)
            if len(layout.placeholders) > 1:
                return layout
        
        # Fallback: Search for a suitable layout
        for layout in prs.slide_layouts:
            if needs_body and len(layout.placeholders) > 1:
                return layout
        
        # Last resort: just return first layout
        return prs.slide_layouts[0]

    # 1. Main Title Slide
    # layout 0 usually title
    title_layout = get_layout(prs, 0) 
    slide = prs.slides.add_slide(title_layout)
    
    if slide.shapes.title:
        title_text = json_data.get("title", "Bài thuyết trình AI")
        slide.shapes.title.text = title_text.upper()
    
    # Safely set subtitle if placeholder exists
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "Được tạo bởi SlideGenius"
        except (IndexError, KeyError):
            pass # Skip if no placeholder accessible

    # 2. Content Slides
    slides_content = json_data.get("slides", [])
    
    # Try to find a content layout (often index 1, needs body)
    content_layout = get_layout(prs, 1, needs_body=True)
    
    # ... existing code ...

    for slide_data in slides_content:
        slide = prs.slides.add_slide(content_layout)
        
        # Set Title and Clean "Slide X:" prefix
        title_height = Cm(0) # Default if no title exists
        if slide.shapes.title:
            raw_title = slide_data.get("title", "")
            # Remove "Slide 1:", "Slide 01", etc.
            clean_title = re.sub(r'^Slide\s+\d+[:.]?\s*', '', raw_title, flags=re.IGNORECASE)
            slide.shapes.title.text = clean_title
            
            # --- Layout Refinement: Advanced Dynamic Calculation ---
            # Instead of a hard threshold, we calculate expected height based on:
            # - Text Length
            # - Font Size (Inherited or Default)
            # - Container Width
            
            # 1. Determine Font Size (Points)
            font_size_pt = 36 # Default
            # Try to read from template shape style if available to be smarter
            # (Limitation: python-pptx often returns None for inherited styles, so we default safe)
            if slide.shapes.title.text_frame and slide.shapes.title.text_frame.paragraphs:
                 p_font = slide.shapes.title.text_frame.paragraphs[0].font
                 if p_font and p_font.size:
                      font_size_pt = p_font.size.pt
            
            # --- FIX: Set Width/Pos BEFORE calculation so math matches reality ---
            # User wants "triệt để", so let's enforce standard margins to be safe.
            slide.shapes.title.left = Cm(1.0)
            slide.shapes.title.width = prs.slide_width - Cm(2.0)
            slide.shapes.title.top = Cm(0.5)

            # 2. Determine Width (Points)
            shape_width_pt = slide.shapes.title.width.pt
            
            # 3. Calculate Capacity
            # Avg char width approx 0.55 of font size for mixed case variable width sans-serif
            avg_char_width = font_size_pt * 0.55
            chars_per_line = shape_width_pt / avg_char_width
            
            # 4. Estimate Lines
            import math
            # Add small buffer +1 for safety
            estimated_lines = math.ceil(len(clean_title) / chars_per_line)
            estimated_lines = max(1, estimated_lines) # At least 1 line
            
            # 5. Calculate Height
            # LINE HEIGHT TUNING:
            # User reported "empty line" effect. Reducing multiplier and removing padding.
            # 1.1 is tight leading.
            line_height_pt = font_size_pt * 1.1
            # Remove extra padding (+0 instead of +12)
            total_height_pt = (estimated_lines * line_height_pt)
            
            title_height = Pt(total_height_pt)
            
            # Apply Height
            slide.shapes.title.height = title_height
            
            # Title Font Styling
            if slide.shapes.title.text_frame:
                tf_title = slide.shapes.title.text_frame
                tf_title.word_wrap = True 
                tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
                tf_title.vertical_anchor = MSO_ANCHOR.TOP # Critical
                
                if tf_title.paragraphs:
                    p = tf_title.paragraphs[0]
                    p.font.size = Pt(font_size_pt)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.LEFT 
                    
                    # Remove margins to make math accurate and remove "empty line" look
                    tf_title.margin_top = 0
                    tf_title.margin_bottom = 0 # Critical fix for "empty line" below
                         
        # Set Content (Body)
        # Find the best placeholder for content
        body_shape = None
        
        # Priority 1: explicitly identified "BODY" or "OBJECT" placeholders
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        candidates = []
        for shape in slide.placeholders:
            # Skip Title, Center Title, Subtitle
            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE]:
                continue
                
            # Prefer Body (2) or Object (7)
            if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                candidates.append((0, shape)) # High priority
            elif shape.placeholder_format.idx == 1:
                candidates.append((1, shape)) # Medium priority
            else:
                candidates.append((2, shape)) # Low priority

        # Sort by priority (asc) then by top/left position implicitly? No, just priority.
        if candidates:
            candidates.sort(key=lambda x: x[0])
            body_shape = candidates[0][1]
        
        # Fallback: if absolutely no specific body, just grab the second placeholder if it exists
        if not body_shape and len(slide.placeholders) > 1:
             # Just grab the one that isn't the title (assuming title is [0] or type title)
             for shape in slide.placeholders:
                 if shape.element is not slide.shapes.title.element:
                      body_shape = shape
                      break

        if body_shape and hasattr(body_shape, "text_frame"):
            # --- Layout Refinement: Adjust Margins & Top ---
            
            margin_side = Cm(1.5)
            # Body Top = Title Top (0.5) + Title Height + Larger Gap (0.5)
            # Dynamic Calculation
            # Balanced gap: 0.4cm (Enough separation, but connected)
            gap = Cm(0.4) 
            margin_top = Cm(0.5) + title_height + gap
            
            # Ensure we don't push off bottom
            margin_bottom = Cm(1.0)
            
            body_shape.left = margin_side
            body_shape.top = margin_top
            body_shape.width = prs.slide_width - (margin_side * 2)
            # Safe calculation for height
            available_height = prs.slide_height - margin_top - margin_bottom
            if available_height < Cm(5): # minimal safe height
                 # If title is huge, shrink body drastically? 
                 # Or just clamp.
                 pass
            
            body_shape.height = available_height
            
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            # --- Hybrid Strategy: Manual Calc + Native Auto-Fit ---
            # 1. We calculate a safe "base" font size manually so the slide looks good immediately on open.
            # 2. We set TEXT_TO_FIT_SHAPE so PowerPoint handles future edits/overflows gracefully.
            
            # Constants for Heuristic
            MAX_FONT_SIZE = 24
            MIN_FONT_SIZE = 10
            
            # Estimate Capacity:
            # Box dimensions: ~22cm wide x ~8cm high.
            # At 24pt, line height ~30pt. Height 8cm=226pt -> ~7 lines.
            # Width 22cm=623pt. Avg char width ~11pt -> ~56 chars/line.
            # Total Chars @ 24pt = 7 * 56 = ~392 chars.
            # Let's be conservative: 350 chars for 24pt.
            # UPDATED: User reported overflow. Drastically reducing logical capacity for safer scaling.
            BASE_CAPACITY_AT_24PT = 300 
            
            total_text_len = sum(len(str(c)) for c in slide_data.get("content", []))
            
            if total_text_len <= BASE_CAPACITY_AT_24PT:
                font_size_pt = MAX_FONT_SIZE
            else:
                # Scaling formula
                import math
                ratio = BASE_CAPACITY_AT_24PT / total_text_len
                scale_factor = math.sqrt(ratio)
                font_size_pt = max(min(MAX_FONT_SIZE * scale_factor, MAX_FONT_SIZE), MIN_FONT_SIZE)
            
            # Apply calculated font size AND set AutoFit
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Re-enable Native Auto-Fit
            
            # Clear default paragraph
            tf.clear()
            
            content = slide_data.get("content", [])
            for i, item in enumerate(content):
                # Reuse the first paragraph if it exists
                if i == 0 and len(tf.paragraphs) == 1:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = 0
                
                # Parse markdown
                parts = re.split(r'(\*\*.*?\*\*)', str(item))
                
                for part in parts:
                    if not part: continue
                    
                    run = p.add_run()
                    # Check for bold marker
                    if part.startswith('**') and part.endswith('**'):
                        text_content = part[2:-2] # Strip **
                        run.text = text_content
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 112, 192) # Emphasis Blue
                    else:
                        run.text = part
                        
                    run.font.size = Pt(font_size_pt) # Set the safe start size
                
                # Adjust spacing based on size
                # Smaller text needs less spacing
                spacing = max(3, font_size_pt * 0.3)
                p.space_before = Pt(spacing)
                p.space_after = Pt(spacing)
        
        # Set Speaker Notes
        
        # Set Speaker Notes
        notes = slide_data.get("notes", "")
        if notes:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
